from docx import Document
from openpyxl import Workbook
from pptx import Presentation

from signalrank.components.data_ingestion.loaders.office import (
    load_office,
)


def test_loads_docx(tmp_path):
    file_path = tmp_path / "sample.docx"

    document = Document()
    document.add_heading("SignalRank Report", level=1)
    document.add_paragraph("Evidence retrieval matters.")
    document.save(file_path)

    elements = load_office(file_path)

    text = "\n".join(
        element.text
        for element in elements
    )

    assert "SignalRank Report" in text
    assert "Evidence retrieval matters." in text


def test_loads_pptx(tmp_path):
    file_path = tmp_path / "sample.pptx"

    presentation = Presentation()

    slide = presentation.slides.add_slide(
        presentation.slide_layouts[1]
    )

    slide.shapes.title.text = "SignalRank"
    slide.placeholders[1].text = "Retrieval evidence"

    presentation.save(file_path)

    elements = load_office(file_path)

    text = "\n".join(
        element.text
        for element in elements
    )

    assert "SignalRank" in text
    assert "Retrieval evidence" in text


def test_loads_xlsx(tmp_path):
    file_path = tmp_path / "sample.xlsx"

    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Evidence"

    sheet["A1"] = "Claim"
    sheet["B1"] = "Support"
    sheet["A2"] = "Drug improves outcome"
    sheet["B2"] = "Study 1"

    workbook.save(file_path)

    elements = load_office(file_path)

    text = "\n".join(
        element.text
        for element in elements
    )

    assert "Claim" in text
    assert "Drug improves outcome" in text